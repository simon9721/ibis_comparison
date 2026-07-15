from __future__ import annotations

import argparse
import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULT_ROOT = ROOT / "results" / "io_buf_two_state_gate_model_2026-06-30"
DEFAULT_TEMPLATE = Path(r"\\minerfiles.mst.edu\dfs\users\sh3qm\Downloads\0710_Simon_IBIS.pptx")
DEFAULT_OUTPUT = RESULT_ROOT / "presentation" / "0714_two_state_directional_residual_presentation.pptx"

GREEN = RGBColor(43, 122, 67)
LIGHT_GREEN = RGBColor(226, 240, 229)
PALE_GREEN = RGBColor(241, 248, 242)
DARK = RGBColor(25, 31, 28)
GRAY = RGBColor(102, 108, 104)
LIGHT_GRAY = RGBColor(242, 243, 242)
MID_GRAY = RGBColor(203, 208, 204)
RED = RGBColor(202, 48, 48)
LIGHT_RED = RGBColor(252, 235, 235)
ORANGE = RGBColor(218, 132, 38)
BLUE = RGBColor(36, 103, 173)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

TITLE_FONT = "Times New Roman"
BODY_FONT = "Aptos"
CODE_FONT = "Consolas"


def add_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    title_shape = slide.shapes.title
    title_shape.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(27)
    p.font.bold = False
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT
    return slide


def move_last_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    slide_id = slide_ids[-1]
    slide_ids.remove(slide_id)
    slide_ids.insert(index, slide_id)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: RGBColor = DARK,
    bold: bool = False,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 17,
    color: RGBColor = DARK,
    spacing: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def add_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor = LIGHT_GREEN,
    line: RGBColor = GREEN,
    size: float = 16,
    bold: bool = False,
    color: RGBColor = DARK,
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = GREEN, width: float = 2.2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_takeaway(slide, text: str, y: float = 6.63, fill: RGBColor = PALE_GREEN, color: RGBColor = GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(y), Inches(12.15), Inches(0.43))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def add_source(slide, text: str):
    return add_text(slide, text, 0.52, 7.16, 11.8, 0.18, size=7.5, color=GRAY)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, border: bool = True):
    with Image.open(path) as image:
        aspect = image.width / image.height
    box_aspect = w / h
    if aspect >= box_aspect:
        pic_w = w
        pic_h = w / aspect
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * aspect
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    if border:
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pic_x - 0.03), Inches(pic_y - 0.03), Inches(pic_w + 0.06), Inches(pic_h + 0.06))
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = MID_GRAY
        frame.line.width = Pt(0.8)
    return slide.shapes.add_picture(str(path), Inches(pic_x), Inches(pic_y), width=Inches(pic_w), height=Inches(pic_h))


def add_code_box(slide, code: str, x: float, y: float, w: float, h: float, size: float = 11.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(246, 248, 247)
    shape.line.color.rgb = RGBColor(164, 174, 168)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0.13)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = CODE_FONT
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(30, 48, 39)
    return shape


def add_metric_card(
    slide,
    label: str,
    value: str,
    status: str,
    x: float,
    y: float,
    w: float,
    status_color: RGBColor,
    value_size: float = 14,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.25))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(1)
    add_text(slide, label, x + 0.12, y + 0.08, w - 0.24, 0.28, size=13.5, color=GRAY, bold=True)
    add_text(slide, value, x + 0.12, y + 0.36, w - 0.24, 0.52, size=value_size, color=DARK, bold=True)
    add_text(slide, status, x + 0.12, y + 0.88, w - 0.24, 0.25, size=12.5, color=status_color, bold=True)
    return card


def add_notes(slide, text: str) -> None:
    """Replace the PowerPoint speaker notes for one slide."""
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.paragraphs[0].text = text.strip()


def add_curve(
    slide,
    points: list[tuple[float, float]],
    x: float,
    y: float,
    w: float,
    h: float,
    color: RGBColor = GREEN,
    width: float = 2.4,
) -> None:
    """Draw a normalized 0..1 curve with native PowerPoint line segments."""
    for (x0, y0), (x1, y1) in zip(points, points[1:]):
        segment = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x + x0 * w),
            Inches(y + (1.0 - y0) * h),
            Inches(x + x1 * w),
            Inches(y + (1.0 - y1) * h),
        )
        segment.line.color.rgb = color
        segment.line.width = Pt(width)


def make_fit_diagnostic_crops() -> tuple[Path, Path]:
    """Split the dense four-panel diagnostic into two teaching-friendly assets."""
    source = RESULT_ROOT / "fit_diagnostics" / "directional_maps_and_residual.png"
    asset_dir = RESULT_ROOT / "presentation" / "assets"
    asset_dir.mkdir(parents=True, exist_ok=True)
    maps_path = asset_dir / "directional_maps_only.png"
    residual_path = asset_dir / "rate_residual_only.png"
    with Image.open(source) as image:
        width, height = image.size
        split = height // 2
        image.crop((0, 0, width, split + 55)).save(maps_path)
        image.crop((0, split - 45, width, height)).save(residual_path)
    return maps_path, residual_path


def build_primer(prs: Presentation) -> None:
    slide = add_slide(prs, "What are Ku and Kd?")
    add_text(slide, "An output buffer is two controlled current networks sharing one pad.", 0.55, 0.93, 12.0, 0.42, size=19, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, "Input\ncommand", 0.65, 2.35, 1.35, 0.85, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_arrow(slide, 2.0, 2.78, 2.75, 2.78)
    add_box(slide, "Output buffer", 2.75, 1.65, 4.25, 2.4, fill=PALE_GREEN, line=GREEN, bold=True, size=18)
    add_box(slide, "Pullup network\nconnects toward 3.3 V", 3.12, 2.15, 3.5, 0.64, fill=LIGHT_GREEN, line=GREEN, size=14)
    add_box(slide, "Pulldown network\nconnects toward ground", 3.12, 3.05, 3.5, 0.64, fill=LIGHT_GREEN, line=GREEN, size=14)
    add_arrow(slide, 7.0, 2.78, 7.8, 2.78)
    add_box(slide, "PAD\nvoltage", 7.8, 2.35, 1.35, 0.85, fill=LIGHT_GRAY, line=GRAY, bold=True)

    add_box(slide, "Ku = pullup strength\n1: fully enabled\n0: off", 9.65, 1.58, 2.75, 1.25, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Kd = pulldown strength\n1: fully enabled\n0: off", 9.65, 3.05, 2.75, 1.25, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)

    add_box(slide, "Input goes HIGH", 1.0, 5.15, 2.1, 0.55, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_arrow(slide, 3.1, 5.43, 3.7, 5.43)
    add_box(slide, "Ku rises; Kd falls", 3.7, 5.15, 2.4, 0.55, fill=LIGHT_GREEN, line=GREEN, bold=True)
    add_arrow(slide, 6.1, 5.43, 6.7, 5.43)
    add_box(slide, "PAD rises", 6.7, 5.15, 1.8, 0.55, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_text(slide, "A short pulse reverses the command before these strengths finish moving.", 8.85, 5.02, 3.5, 0.85, size=15, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Ku/Kd describe effective drive strength. The hard case is preserving their unfinished history.")
    add_source(slide, "Conceptual diagram; Ku/Kd interpretation follows the io_buf native-IBIS and pybis coefficient convention.")
    move_last_slide(prs, 1)


def build_setup_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Controlled experiment: same stimulus, three references")
    add_text(slide, "Testbench", 0.62, 1.0, 2.0, 0.35, size=19, bold=True, color=GREEN)
    add_box(slide, "1 ps edge\n3.3 V input", 0.7, 1.55, 1.75, 0.9, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_arrow(slide, 2.45, 2.0, 3.2, 2.0)
    add_box(slide, "io_buf driver", 3.2, 1.55, 2.0, 0.9, fill=LIGHT_GREEN, line=GREEN, bold=True)
    add_arrow(slide, 5.2, 2.0, 5.95, 2.0)
    add_box(slide, "PAD", 5.95, 1.68, 1.2, 0.64, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_arrow(slide, 7.15, 2.0, 7.9, 2.0)
    add_box(slide, "50 ohm || 2 pF", 7.9, 1.55, 2.1, 0.9, fill=LIGHT_GRAY, line=GRAY, bold=True)

    add_text(slide, "What each flow tells us", 0.62, 3.0, 3.2, 0.35, size=19, bold=True, color=GREEN)
    add_box(slide, "HSPICE native IBIS\nPad + Ku + Kd\nPlayback reference", 0.72, 3.55, 3.65, 1.25, fill=RGBColor(245, 245, 245), line=BLACK, size=16, bold=True)
    add_box(slide, "HSPICE transistor io_buf.sp\nPad only\nCircuit-level reference", 4.82, 3.55, 3.65, 1.25, fill=LIGHT_GRAY, line=GRAY, size=16, bold=True)
    add_box(slide, "ngspice + pybis candidate\nPad + Ku + Kd + hidden states\nAlgorithm under test", 8.92, 3.55, 3.65, 1.25, fill=LIGHT_GREEN, line=GREEN, size=16, bold=True)

    add_box(slide, "Long pulse", 1.0, 5.4, 2.1, 0.58, fill=PALE_GREEN, line=GREEN, bold=True)
    add_text(slide, "Normal-operation regression guard", 3.25, 5.45, 3.4, 0.36, size=16)
    add_box(slide, "Short high / short low", 7.0, 5.4, 2.6, 0.58, fill=LIGHT_RED, line=RED, bold=True)
    add_text(slide, "Interrupted-transition challenge", 9.75, 5.45, 2.65, 0.36, size=16)
    add_takeaway(slide, "HSPICE validates the result; no HSPICE waveform is used to fit the ngspice model.")
    add_source(slide, "Cached study setup: results/io_buf_two_state_gate_model_2026-06-30")


def build_architecture_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Proposed solution: remember two hidden gate states")
    stages = [
        ("Input edge", 0.45, LIGHT_GRAY, GRAY),
        ("Directional\ndelays", 2.18, LIGHT_GREEN, GREEN),
        ("GUP / GDN\ncontinuous states", 4.02, LIGHT_GREEN, GREEN),
        ("Direction-specific\nPWL maps", 6.08, LIGHT_GREEN, GREEN),
        ("Kd rate\nresidual", 8.23, RGBColor(255, 244, 224), ORANGE),
        ("Ku / Kd", 10.12, RGBColor(234, 244, 255), BLUE),
        ("IBIS currents\nand PAD", 11.48, LIGHT_GRAY, GRAY),
    ]
    widths = [1.35, 1.45, 1.7, 1.75, 1.45, 1.05, 1.35]
    for idx, ((label, x, fill, line), w) in enumerate(zip(stages, widths)):
        add_box(slide, label, x, 2.05, w, 1.05, fill=fill, line=line, size=14, bold=True)
        if idx < len(stages) - 1:
            next_x = stages[idx + 1][1]
            add_arrow(slide, x + w, 2.58, next_x, 2.58, color=GREEN, width=1.8)

    add_text(slide, "What changes at a reverse edge?", 0.75, 3.75, 4.2, 0.35, size=19, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "The command target changes direction.",
            "GUP and GDN continue from their present values.",
            "Ku and Kd remain continuous; nothing restarts at t = 0.",
        ],
        0.8,
        4.18,
        5.4,
        1.72,
        size=16,
    )
    add_text(slide, "What is learned from the IBIS tables?", 6.8, 3.75, 5.3, 0.35, size=19, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "Four onset delays and four time constants.",
            "Separate on/off mappings for pullup and pulldown.",
            "A small Kd correction linked to gate-state rate.",
        ],
        6.85,
        4.18,
        5.4,
        1.72,
        size=16,
    )
    add_takeaway(slide, "Best structural method so far: two_state_directional_residual.")
    add_source(slide, "Implementation: tools/pybis2spice/pybis2spice/subcircuit.py")


def build_state_memory_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "How a capacitor stores unfinished switching history")
    add_text(slide, "Pullup state", 0.7, 1.05, 2.1, 0.35, size=19, bold=True, color=BLUE)
    add_box(slide, "GUPTARGET", 0.75, 1.65, 1.65, 0.65, fill=RGBColor(234, 244, 255), line=BLUE, bold=True)
    add_arrow(slide, 2.4, 1.98, 3.2, 1.98, color=BLUE)
    add_box(slide, "C_GUP\nstored state GUP", 3.2, 1.42, 2.25, 1.12, fill=WHITE, line=BLUE, bold=True)
    add_text(slide, "dGUP/dt = (target - GUP) / tau_PU", 0.82, 2.82, 4.8, 0.48, size=16, font=CODE_FONT, color=BLUE, bold=True)

    add_text(slide, "Pulldown state", 6.85, 1.05, 2.2, 0.35, size=19, bold=True, color=ORANGE)
    add_box(slide, "GDNTARGET", 6.9, 1.65, 1.65, 0.65, fill=RGBColor(255, 241, 228), line=ORANGE, bold=True)
    add_arrow(slide, 8.55, 1.98, 9.35, 1.98, color=ORANGE)
    add_box(slide, "C_GDN\nstored state GDN", 9.35, 1.42, 2.25, 1.12, fill=WHITE, line=ORANGE, bold=True)
    add_text(slide, "dGDN/dt = (target - GDN) / tau_PD", 6.98, 2.82, 4.8, 0.48, size=16, font=CODE_FONT, color=ORANGE, bold=True)

    add_box(slide, "Before reverse edge\nGUP = 0.31\nGDN = 0.42", 0.85, 4.15, 2.15, 1.2, fill=LIGHT_GRAY, line=GRAY, size=16, bold=True)
    add_arrow(slide, 3.0, 4.75, 4.15, 4.75, color=RED)
    add_box(slide, "Reverse edge\nchanges targets", 4.15, 4.28, 2.05, 0.92, fill=LIGHT_RED, line=RED, size=16, bold=True)
    add_arrow(slide, 6.2, 4.75, 7.35, 4.75, color=GREEN)
    add_box(slide, "After reverse edge\nGUP starts at 0.31\nGDN starts at 0.42", 7.35, 4.15, 2.5, 1.2, fill=LIGHT_GREEN, line=GREEN, size=16, bold=True)
    add_text(slide, "State is preserved", 10.15, 4.48, 2.0, 0.45, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    add_takeaway(slide, "The capacitor is the memory: voltage cannot jump, so a reverse edge continues from the true partial state.")
    add_source(slide, "Conceptual state equation; generated ngspice implementation uses behavioral current sources plus capacitors.")


def build_direction_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Direction matters: four processes, four fitted dynamics")
    add_text(slide, "Input rises", 0.72, 1.0, 5.8, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Input falls", 6.8, 1.0, 5.8, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_box(slide, "Pullup turns ON\nPU delay 1.302 ns\nPU tau 1.284 ns\nuse Ku-on map", 0.82, 1.58, 2.55, 1.45, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Pulldown turns OFF\nPD delay 1.443 ns\nPD tau 0.237 ns\nuse Kd-off map", 3.68, 1.58, 2.55, 1.45, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)
    add_box(slide, "Pullup turns OFF\nPU delay 0.440 ns\nPU tau 0.363 ns\nuse Ku-off map", 6.9, 1.58, 2.55, 1.45, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Pulldown turns ON\nPD delay 2.480 ns\nPD tau 0.275 ns\nuse Kd-on map", 9.76, 1.58, 2.55, 1.45, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)

    add_text(slide, "Why not one shared map?", 0.78, 3.62, 3.3, 0.38, size=20, bold=True, color=RED)
    add_bullets(
        slide,
        [
            "The same gate-state value can produce a different coefficient while turning on versus turning off.",
            "Pullup and pulldown have different predriver delays and device physics.",
            "A short pulse can interrupt any one of these four paths.",
        ],
        0.82,
        4.05,
        6.0,
        1.72,
        size=16,
    )
    add_box(slide, "Decision at runtime", 7.35, 3.68, 2.1, 0.48, fill=LIGHT_GRAY, line=GRAY, bold=True)
    add_arrow(slide, 9.45, 3.92, 10.15, 3.92)
    add_box(slide, "target >= state?\nuse ON map\nelse use OFF map", 10.15, 3.42, 2.05, 1.02, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_box(slide, "Kd only: add a small rate residual to reproduce undershoot missed by static maps.", 7.35, 4.95, 4.85, 0.9, fill=RGBColor(255, 244, 224), line=ORANGE, size=16, bold=True)
    add_takeaway(slide, "Direction-specific maps fixed the large shape error; the Kd rate residual recovered the missing undershoot.")
    add_source(slide, "Numbers are generated from io_buf IBIS-derived coefficient tables; see driver_OutputInput_Typical.sub comments.")


def build_training_data_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 0: the model learns from four complete-edge traces")
    add_text(slide, "Rising input command", 0.72, 0.98, 5.65, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Falling input command", 6.9, 0.98, 5.65, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_box(slide, "kr[:, TIME]", 0.82, 1.55, 1.55, 0.58, fill=LIGHT_GRAY, line=GRAY, size=14, bold=True)
    add_box(slide, "kr[:, Ku]\npullup turns ON", 2.62, 1.42, 1.95, 0.84, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_box(slide, "kr[:, Kd]\npulldown turns OFF", 4.82, 1.42, 1.95, 0.84, fill=RGBColor(255, 241, 228), line=ORANGE, size=14, bold=True)
    add_box(slide, "kf[:, TIME]", 7.0, 1.55, 1.55, 0.58, fill=LIGHT_GRAY, line=GRAY, size=14, bold=True)
    add_box(slide, "kf[:, Ku]\npullup turns OFF", 8.8, 1.42, 1.95, 0.84, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_box(slide, "kf[:, Kd]\npulldown turns ON", 11.0, 1.42, 1.95, 0.84, fill=RGBColor(255, 241, 228), line=ORANGE, size=14, bold=True)

    add_text(slide, "First infer the settled endpoints", 0.78, 2.88, 4.5, 0.38, size=20, bold=True, color=GREEN)
    add_code_box(
        slide,
        "Ku_off = mean(Ku_rise[0],  Ku_fall[-1])\n"
        "Ku_on  = mean(Ku_rise[-1], Ku_fall[0])\n"
        "Kd_on  = mean(Kd_rise[0],  Kd_fall[-1])\n"
        "Kd_off = mean(Kd_rise[-1], Kd_fall[0])",
        0.78,
        3.35,
        6.0,
        1.72,
        size=13.3,
    )
    add_text(slide, "io_buf endpoints used by this run", 7.2, 2.88, 4.9, 0.38, size=20, bold=True, color=GREEN)
    add_box(slide, "Ku_off = 0.00194\nKu_on  = 0.99579", 7.22, 3.45, 2.45, 1.08, fill=RGBColor(234, 244, 255), line=BLUE, size=16, bold=True)
    add_box(slide, "Kd_off = 0.00110\nKd_on  = 0.99861", 9.98, 3.45, 2.45, 1.08, fill=RGBColor(255, 241, 228), line=ORANGE, size=16, bold=True)

    add_box(slide, "Training input", 0.82, 5.58, 1.65, 0.5, fill=LIGHT_GRAY, line=GRAY, size=14, bold=True)
    add_arrow(slide, 2.47, 5.83, 3.18, 5.83)
    add_box(slide, "IBIS-derived Ku/Kd tables only", 3.18, 5.45, 3.55, 0.76, fill=LIGHT_GREEN, line=GREEN, size=16, bold=True)
    add_arrow(slide, 6.73, 5.83, 7.45, 5.83)
    add_box(slide, "No HSPICE waveform enters the fit", 7.45, 5.45, 4.3, 0.76, fill=LIGHT_RED, line=RED, size=16, bold=True)
    add_takeaway(slide, "Everything that follows is derived from these four complete-edge coefficient traces.")
    add_source(slide, "Implementation: gate_state_fit(kr, kf) in tools/pybis2spice/pybis2spice/subcircuit.py")


def build_progress_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 1: normalize each coefficient into transition progress")
    add_text(slide, "One formula handles both increasing and decreasing traces", 0.62, 0.98, 12.1, 0.4, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_code_box(slide, "progress(t) = [K(t) - K_start] / [K_end - K_start]", 1.8, 1.55, 9.7, 0.72, size=18)

    labels = [("start", "0%"), ("t05", "5%"), ("t63", "63.2%"), ("t90", "90%"), ("end", "100%")]
    xs = [0.9, 3.1, 6.0, 9.0, 11.25]
    for idx, ((head, value), xpos) in enumerate(zip(labels, xs)):
        fill = LIGHT_GREEN if idx in (1, 2, 3) else LIGHT_GRAY
        line = GREEN if idx in (1, 2, 3) else GRAY
        add_box(slide, f"{head}\n{value}", xpos, 3.0, 1.25, 0.88, fill=fill, line=line, size=15, bold=True)
        if idx < len(xs) - 1:
            add_arrow(slide, xpos + 1.25, 3.44, xs[idx + 1], 3.44, color=GREEN, width=1.7)

    add_text(slide, "Why normalize?", 0.78, 4.55, 2.4, 0.35, size=19, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "Ku and Kd have different numerical endpoints, but progress always runs from 0 to 1.",
            "The sign of K_end - K_start automatically handles a falling coefficient.",
            "The code keeps a wider -0.5 to 1.5 range so real overshoot or undershoot is not erased.",
        ],
        0.82,
        4.95,
        11.5,
        1.35,
        size=15.5,
        spacing=5,
    )
    add_takeaway(slide, "Normalized progress lets the same timing extractor operate on Ku-on, Ku-off, Kd-on, and Kd-off.")
    add_source(slide, "Actual helpers: coefficient_progress() and crossing_time_ns()")


def build_tau_derivation_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 2: derive delay and tau from 5%, 63.2%, and 90% crossings")

    plot_x, plot_y, plot_w, plot_h = 0.72, 1.45, 5.35, 3.35
    x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(plot_x), Inches(plot_y + plot_h), Inches(plot_x + plot_w), Inches(plot_y + plot_h))
    y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(plot_x), Inches(plot_y + plot_h), Inches(plot_x), Inches(plot_y))
    for axis in (x_axis, y_axis):
        axis.line.color.rgb = GRAY
        axis.line.width = Pt(1.4)
    delay_fraction = 0.2
    tau_fraction = 0.22
    points = []
    for idx in range(61):
        xn = idx / 60
        if xn <= delay_fraction:
            yn = 0.0
        else:
            yn = 1.0 - math.exp(-(xn - delay_fraction) / tau_fraction)
        points.append((xn, min(yn, 1.0)))
    add_curve(slide, points, plot_x, plot_y, plot_w, plot_h, color=GREEN, width=2.8)
    for level, label, xpos in [(0.05, "t05", 0.211), (0.632, "t63", 0.42), (0.90, "t90", 0.705)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(plot_x), Inches(plot_y + (1 - level) * plot_h), Inches(plot_x + xpos * plot_w), Inches(plot_y + (1 - level) * plot_h))
        line.line.color.rgb = MID_GRAY
        line.line.width = Pt(1)
        add_text(slide, label, plot_x + xpos * plot_w - 0.22, plot_y + (1 - level) * plot_h - 0.34, 0.6, 0.24, size=11, color=GRAY, bold=True)
    add_text(slide, "normalized progress", 0.05, 2.55, 1.55, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "table time", 2.75, 4.84, 1.4, 0.28, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, "Exact extraction rule used in code", 6.48, 1.15, 5.9, 0.38, size=20, bold=True, color=GREEN)
    add_code_box(
        slide,
        "delay = max(0, t05)\n"
        "tau63 = t63 - delay\n"
        "tau90 = (t90 - delay) / ln(10)\n"
        "tau = max(0.020 ns, tau63, tau90)",
        6.48,
        1.64,
        5.75,
        1.82,
        size=15,
    )
    add_box(slide, "Why 63.2%?\nA first-order response reaches 1 - e^-1 = 63.2% after one tau.", 6.48, 3.8, 2.72, 1.22, fill=LIGHT_GREEN, line=GREEN, size=14, bold=True)
    add_box(slide, "Why also 90%?\nIt prevents a fast early crossing from underestimating the slower tail.", 9.5, 3.8, 2.72, 1.22, fill=RGBColor(255, 244, 224), line=ORANGE, size=14, bold=True)
    add_text(slide, "This is a deterministic timing estimate, not a nonlinear optimizer.", 6.6, 5.38, 5.5, 0.4, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Delay captures when motion starts; tau captures how quickly the hidden state approaches its target.")
    add_source(slide, "Actual helper: coefficient_transition_timing(); all time values are converted to ns before fitting.")


def build_hidden_state_reconstruction_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 3: reconstruct GUP and GDN from the fitted timing")
    add_text(slide, "The hidden state is normalized internal progress, not the final Ku or Kd coefficient.", 0.62, 0.98, 12.1, 0.38, size=19, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_code_box(
        slide,
        "x = max(0, t - delay)\n"
        "progress = 1 - exp(-x / tau)\n"
        "G(t) = G_start + (G_end - G_start) * progress",
        1.2,
        1.55,
        10.9,
        1.45,
        size=17,
    )

    add_box(slide, "Input rises\nGUP: 0 -> 1\nPU-on delay/tau", 0.72, 3.55, 2.55, 1.2, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Input rises\nGDN: 1 -> 0\nPD-off delay/tau", 3.48, 3.55, 2.55, 1.2, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)
    add_box(slide, "Input falls\nGUP: 1 -> 0\nPU-off delay/tau", 6.24, 3.55, 2.55, 1.2, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Input falls\nGDN: 0 -> 1\nPD-on delay/tau", 9.0, 3.55, 2.55, 1.2, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)

    add_box(slide, "Initial low output: GUP = 0", 1.2, 5.35, 3.3, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_box(slide, "Initial low output: GDN = 1", 4.95, 5.35, 3.3, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_box(slide, "Reverse edge changes target; G itself is preserved", 8.7, 5.25, 3.5, 0.82, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_takeaway(slide, "At every original table sample time, the fitted delay/tau gives one corresponding hidden-state value.")
    add_source(slide, "Actual helper: gate_response(time_ns, delay_ns, tau_ns, start_value, end_value)")


def build_pwl_derivation_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 4: turn hidden state and coefficient samples into four PWL maps")
    flow = [
        ("For each table time t", LIGHT_GRAY, GRAY),
        ("Compute fitted G(t)", LIGHT_GREEN, GREEN),
        ("Pair (G(t), K(t))", LIGHT_GREEN, GREEN),
        ("Sort by G and merge duplicates", LIGHT_GREEN, GREEN),
        ("Interpolate 81 points from G=0..1", RGBColor(255, 244, 224), ORANGE),
    ]
    xpos = [0.45, 2.85, 5.2, 7.55, 9.9]
    widths = [1.9, 1.85, 1.85, 1.9, 2.45]
    for idx, ((label, fill, line), x, w) in enumerate(zip(flow, xpos, widths)):
        add_box(slide, label, x, 1.42, w, 0.92, fill=fill, line=line, size=13.5, bold=True)
        if idx < len(flow) - 1:
            add_arrow(slide, x + w, 1.88, xpos[idx + 1], 1.88, color=GREEN, width=1.5)

    add_text(slide, "Why four maps?", 0.72, 2.95, 2.4, 0.38, size=20, bold=True, color=GREEN)
    add_box(slide, "Ku-on\nkr Ku vs GUP", 0.75, 3.5, 2.1, 0.92, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Ku-off\nkf Ku vs GUP", 3.05, 3.5, 2.1, 0.92, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_box(slide, "Kd-off\nkr Kd vs GDN", 5.35, 3.5, 2.1, 0.92, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)
    add_box(slide, "Kd-on\nkf Kd vs GDN", 7.65, 3.5, 2.1, 0.92, fill=RGBColor(255, 241, 228), line=ORANGE, size=15, bold=True)
    add_box(slide, "Same G can map to different K while turning on versus turning off", 10.05, 3.32, 2.5, 1.28, fill=LIGHT_RED, line=RED, size=14, bold=True)

    add_text(slide, "Critical implementation choice", 0.72, 5.0, 3.2, 0.35, size=19, bold=True, color=RED)
    add_text(slide, "The directional maps are NOT forced monotonic. That preserves the real negative Kd excursion instead of smoothing it away.", 0.78, 5.44, 11.6, 0.64, size=16.5, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "The PWL map converts state into coefficient; direction selects which history-dependent map is valid.")
    add_source(slide, "Actual helper: directional_gate_transfer_curve(..., point_count=81)")


def build_rate_residual_derivation_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Step 5: recover the Kd detail that a static state map misses")
    add_text(slide, "A single value of GDN does not capture every dynamic feature of Kd, especially its negative undershoot.", 0.62, 0.98, 12.1, 0.4, size=18.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_box(slide, "1. Base prediction", 0.72, 1.62, 2.2, 0.5, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_code_box(slide, "Kd_base(t) = PWL_direction[GDN(t)]", 0.72, 2.2, 3.75, 0.72, size=14.2)
    add_box(slide, "2. Table residual", 4.75, 1.62, 2.2, 0.5, fill=RGBColor(255, 244, 224), line=ORANGE, size=15, bold=True)
    add_code_box(slide, "r(t) = Kd_table(t) - Kd_base(t)", 4.75, 2.2, 3.75, 0.72, size=14.2)
    add_box(slide, "3. Gate-state rate", 8.78, 1.62, 2.2, 0.5, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_code_box(slide, "dGDN/dt = (target - GDN) / tau", 8.78, 2.2, 3.75, 0.72, size=14.2)

    add_text(slide, "Fit one scalar rate gain by least squares", 0.75, 3.42, 5.5, 0.38, size=20, bold=True, color=GREEN)
    add_code_box(
        slide,
        "a = sum[(dGDN/dt) * r] / sum[(dGDN/dt)^2]\n"
        "a = 0.0002386915 ns for io_buf",
        0.78,
        3.92,
        5.65,
        1.22,
        size=15,
    )
    add_text(slide, "Runtime Kd construction", 6.82, 3.42, 3.8, 0.38, size=20, bold=True, color=GREEN)
    add_code_box(
        slide,
        "Kd = Kd_base\n"
        "   + Kd_residual_table(edge_time)\n"
        "   + a * dGDN/dt",
        6.85,
        3.92,
        5.2,
        1.42,
        size=16,
    )
    add_box(slide, "Important nuance: the main state path is continuous, but the small residual-table correction still uses elapsed edge time.", 0.95, 5.62, 11.4, 0.62, fill=LIGHT_RED, line=RED, size=15.5, bold=True)
    add_takeaway(slide, "The rate term adds direction and speed information that GDN amplitude alone cannot express.")
    add_source(slide, "Actual fit: kd_rate_gain_ns = dot(kd_rate, kd_residual) / dot(kd_rate, kd_rate)")


def build_map_diagnostic_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Reading the fit evidence: the on and off maps are genuinely different")
    maps_path, _ = make_fit_diagnostic_crops()
    add_picture_contain(slide, maps_path, 0.5, 1.08, 9.0, 4.85)
    add_box(slide, "Horizontal axis\nnormalized hidden state\nGUP or GDN", 9.8, 1.25, 2.65, 1.0, fill=LIGHT_GRAY, line=GRAY, size=14, bold=True)
    add_box(slide, "Vertical axis\neffective coefficient\nKu or Kd", 9.8, 2.55, 2.65, 1.0, fill=LIGHT_GRAY, line=GRAY, size=14, bold=True)
    add_box(slide, "Blue vs orange separation means one shared map would erase directional history.", 9.8, 3.85, 2.65, 1.45, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_takeaway(slide, "The maps are learned point by point from complete edges; they are not hand-drawn transfer curves.")
    add_source(slide, str(maps_path.relative_to(ROOT)))


def build_residual_diagnostic_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Reading the fit evidence: rate and residual explain the Kd undershoot")
    _, residual_path = make_fit_diagnostic_crops()
    add_picture_contain(slide, residual_path, 0.5, 1.08, 9.0, 4.85)
    add_box(slide, "Left panel\ndGDN/dt is negative when pulldown turns off and positive when it turns on.", 9.8, 1.18, 2.65, 1.3, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_box(slide, "Right panel\nresidual is what the PWL state map did not reproduce.", 9.8, 2.78, 2.65, 1.18, fill=RGBColor(255, 244, 224), line=ORANGE, size=14, bold=True)
    add_box(slide, "The large negative region is why a static nonnegative map failed the reconstruction gate.", 9.8, 4.26, 2.65, 1.18, fill=LIGHT_RED, line=RED, size=14, bold=True)
    add_takeaway(slide, "The residual is small in model structure, but essential for preserving the measured Kd excursion.")
    add_source(slide, str(residual_path.relative_to(ROOT)))


def build_fit_diagnostics_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "The mappings and residual are learned from the IBIS tables")
    path = RESULT_ROOT / "fit_diagnostics" / "directional_maps_and_residual.png"
    add_picture_contain(slide, path, 0.48, 0.96, 12.35, 5.55)
    add_takeaway(slide, "Top: separate on/off mappings. Bottom: Kd residual follows the part a static state map cannot reproduce.")
    add_source(slide, str(path.relative_to(ROOT)))


def build_software_flow_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "End-to-end implementation flow")
    labels = [
        ("io_buf.ibs", "Input data"),
        ("pybis extraction", "Ku/Kd rise + fall tables"),
        ("directional fit", "delays, taus, maps, residual"),
        ("SPICE generator", "state nodes + PWL equations"),
        ("ngspice deck", "50 ohm || 2 pF transient"),
        ("cached audit", "native IBIS + transistor pad"),
    ]
    x_positions = [0.45, 2.48, 4.52, 6.56, 8.6, 10.64]
    for idx, ((head, sub), x) in enumerate(zip(labels, x_positions)):
        fill = LIGHT_GREEN if idx in (2, 3) else LIGHT_GRAY
        line = GREEN if idx in (2, 3) else GRAY
        add_box(slide, f"{head}\n{sub}", x, 2.02, 1.7, 1.18, fill=fill, line=line, size=13.5, bold=True)
        if idx < len(labels) - 1:
            add_arrow(slide, x + 1.7, 2.61, x_positions[idx + 1], 2.61, color=GREEN, width=1.6)

    add_text(slide, "Primary Python entry points", 0.72, 3.85, 3.3, 0.38, size=20, bold=True, color=GREEN)
    add_code_box(
        slide,
        "two_state_directional_gate_fit(kr, kf)\n"
        "create_ngspice_two_state_gate_input_control_netlist(\n"
        "    kr, kf, ibis_data, mode='directional_residual_full')",
        0.75,
        4.28,
        6.05,
        1.35,
        size=13,
    )
    add_text(slide, "Generated artifacts", 7.25, 3.85, 2.8, 0.38, size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "driver_OutputInput_Typical.sub",
            "ngspice transient deck and raw file",
            "diagnostic nodes: GUP, GDN, Ku, Kd, residual",
        ],
        7.3,
        4.25,
        5.0,
        1.55,
        size=15.5,
    )
    add_takeaway(slide, "Normal generation needs only the IBIS file and ngspice. HSPICE enters after model selection as an audit.")
    add_source(slide, "tools/pybis2spice/pybis2spice/subcircuit.py; scripts/run_io_buf_two_state_gate_model.py")


def build_timing_code_detail_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Implementation 1: endpoint and timing extraction in Python")
    code = (
        "ku_off = mean([kr[0, KU],  kf[-1, KU]])\n"
        "ku_on  = mean([kr[-1, KU], kf[0, KU]])\n"
        "kd_on  = mean([kr[0, KD],  kf[-1, KD]])\n"
        "kd_off = mean([kr[-1, KD], kf[0, KD]])\n"
        "\n"
        "def coefficient_transition_timing(time, values, start, end):\n"
        "    time_ns, progress = coefficient_progress(time, values, start, end)\n"
        "    t05 = crossing_time_ns(time_ns, progress, 0.05)\n"
        "    t63 = crossing_time_ns(time_ns, progress, 0.632)\n"
        "    t90 = crossing_time_ns(time_ns, progress, 0.90)\n"
        "    delay = max(0.0, t05)\n"
        "    tau = max(0.02, t63 - delay)\n"
        "    tau = max(tau, (t90 - delay) / 2.302585093)\n"
        "    return delay, tau"
    )
    add_code_box(slide, code, 0.52, 1.02, 8.15, 5.45, size=11.1)
    add_text(slide, "Line-by-line meaning", 8.98, 1.08, 3.25, 0.38, size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "Average both tables at the settled endpoints to reduce endpoint noise.",
            "Normalize each trace so a decreasing Kd is treated exactly like an increasing Ku.",
            "Interpolate the first crossing; do not require an exact sample at 5%, 63.2%, or 90%.",
            "Use the slower of the 63.2% and 90% estimates, with a 20 ps floor.",
        ],
        8.98,
        1.58,
        3.45,
        3.82,
        size=14.3,
        spacing=7,
    )
    add_box(slide, "Output: four delays + four taus", 9.05, 5.55, 3.25, 0.62, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_takeaway(slide, "This stage turns raw coefficient tables into compact, reproducible timing parameters.")
    add_source(slide, "Exact implementation shortened only by replacing internal column constants with TIME/Ku/Kd labels.")


def build_map_fit_code_detail_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Implementation 2: build directional PWL maps and fit the rate gain")
    code = (
        "gdn_rise = gate_response(tr, pd_off_delay, pd_off_tau, 1.0, 0.0)\n"
        "gdn_fall = gate_response(tf, pd_on_delay,  pd_on_tau,  0.0, 1.0)\n"
        "\n"
        "kd_off_x, kd_off_y = directional_gate_transfer_curve(\n"
        "    kr[:, TIME], kr[:, KD], pd_off_delay, pd_off_tau, 1.0, 0.0,\n"
        "    kd_off, kd_on)\n"
        "kd_on_x, kd_on_y = directional_gate_transfer_curve(\n"
        "    kf[:, TIME], kf[:, KD], pd_on_delay, pd_on_tau, 0.0, 1.0,\n"
        "    kd_off, kd_on)\n"
        "\n"
        "kd_base = concatenate([interp(gdn_rise, kd_off_x, kd_off_y),\n"
        "                       interp(gdn_fall, kd_on_x, kd_on_y)])\n"
        "kd_residual = kd_table - kd_base\n"
        "kd_rate_gain = dot(kd_rate, kd_residual) / dot(kd_rate, kd_rate)"
    )
    add_code_box(slide, code, 0.52, 1.02, 8.55, 5.48, size=10.5)
    add_box(slide, "gate_response\ncreates the hidden state at each table time", 9.38, 1.18, 3.0, 0.9, fill=LIGHT_GREEN, line=GREEN, size=14, bold=True)
    add_box(slide, "directional_gate_transfer_curve\npairs state with the original coefficient", 9.38, 2.45, 3.0, 1.02, fill=LIGHT_GREEN, line=GREEN, size=14, bold=True)
    add_box(slide, "interp\nreconstructs Kd from the state-only map", 9.38, 3.85, 3.0, 0.9, fill=RGBColor(255, 244, 224), line=ORANGE, size=14, bold=True)
    add_box(slide, "dot / dot\nfits one least-squares scalar for the rate term", 9.38, 5.12, 3.0, 0.9, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_takeaway(slide, "The fit remains global and deterministic; there is no per-pulse tuning in this code path.")
    add_source(slide, "Actual implementation: two_state_directional_gate_fit(kr, kf)")


def build_command_target_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Implementation 3: each input edge launches four delayed gate commands")
    add_text(slide, "Rising input edge", 0.7, 1.0, 5.8, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Falling input edge", 6.8, 1.0, 5.8, 0.38, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_box(slide, "RISEEDGE", 0.8, 1.62, 1.65, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_arrow(slide, 2.45, 1.93, 3.2, 1.93)
    add_box(slide, "delay 1.302 ns\nPUONP", 3.2, 1.45, 2.0, 0.96, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_arrow(slide, 5.2, 1.93, 5.85, 1.93)
    add_box(slide, "GUPCMD +1\nGUPTARGET -> 1", 5.85, 1.45, 2.0, 0.96, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)

    add_box(slide, "RISEEDGE", 0.8, 3.0, 1.65, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_arrow(slide, 2.45, 3.31, 3.2, 3.31)
    add_box(slide, "delay 1.443 ns\nPDOFFP", 3.2, 2.83, 2.0, 0.96, fill=RGBColor(255, 241, 228), line=ORANGE, size=14, bold=True)
    add_arrow(slide, 5.2, 3.31, 5.85, 3.31)
    add_box(slide, "GDNCMD -1\nGDNTARGET -> 0", 5.85, 2.83, 2.0, 0.96, fill=RGBColor(255, 241, 228), line=ORANGE, size=14, bold=True)

    add_box(slide, "FALLEDGE", 7.1, 1.62, 1.65, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_arrow(slide, 8.75, 1.93, 9.5, 1.93)
    add_box(slide, "delay 0.440 ns\nPUOFFP", 9.5, 1.45, 2.0, 0.96, fill=RGBColor(234, 244, 255), line=BLUE, size=14, bold=True)
    add_arrow(slide, 11.5, 1.93, 12.1, 1.93)
    add_box(slide, "GUPCMD -1", 12.1, 1.57, 0.9, 0.72, fill=RGBColor(234, 244, 255), line=BLUE, size=12.5, bold=True)

    add_box(slide, "FALLEDGE", 7.1, 3.0, 1.65, 0.62, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)
    add_arrow(slide, 8.75, 3.31, 9.5, 3.31)
    add_box(slide, "delay 2.480 ns\nPDONP", 9.5, 2.83, 2.0, 0.96, fill=RGBColor(255, 241, 228), line=ORANGE, size=14, bold=True)
    add_arrow(slide, 11.5, 3.31, 12.1, 3.31)
    add_box(slide, "GDNCMD +1", 12.1, 2.95, 0.9, 0.72, fill=RGBColor(255, 241, 228), line=ORANGE, size=12.5, bold=True)

    add_text(slide, "The delayed pulses update command capacitors; target nodes are clamped to 0..1.", 0.75, 4.55, 12.0, 0.4, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_code_box(
        slide,
        "GUPTARGET = clamp(GUPCMD, 0, 1)\n"
        "GDNTARGET = clamp(GDNCMD, 0, 1)",
        3.55,
        5.12,
        6.25,
        0.96,
        size=16,
    )
    add_takeaway(slide, "A reverse edge schedules the opposite commands; the stored GUP/GDN voltages are not reset.")
    add_source(slide, "Generated SPICE uses lossless T-lines as transport delays for PUONP, PUOFFP, PDOFFP, and PDONP.")


def build_fit_code_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Code demo 1: fit directional maps and the Kd residual")
    code = (
        "def two_state_directional_gate_fit(kr, kf):\n"
        "    fit = dict(gate_state_fit(kr, kf))\n"
        "\n"
        "    ku_on = directional_gate_transfer_curve(\n"
        "        kr[:, TIME], kr[:, KU], fit['pu_on_delay'], fit['pu_on_tau'], 0, 1)\n"
        "    ku_off = directional_gate_transfer_curve(\n"
        "        kf[:, TIME], kf[:, KU], fit['pu_off_delay'], fit['pu_off_tau'], 1, 0)\n"
        "    kd_off = directional_gate_transfer_curve(...)\n"
        "    kd_on  = directional_gate_transfer_curve(...)\n"
        "\n"
        "    kd_base = map_gate_state_to_kd(...)\n"
        "    kd_residual = kd_table - kd_base\n"
        "    gain = dot(dGDN_dt, kd_residual) / dot(dGDN_dt, dGDN_dt)\n"
        "    fit['kd_rate_gain_ns'] = gain\n"
        "    return fit"
    )
    add_code_box(slide, code, 0.55, 1.05, 8.05, 5.35, size=11.4)
    add_box(slide, "1. Start with four delays / taus", 8.9, 1.25, 3.65, 0.72, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_box(slide, "2. Convert complete-edge time into hidden gate state", 8.9, 2.25, 3.65, 0.92, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_box(slide, "3. Pair gate state with measured Ku/Kd to make four PWL maps", 8.9, 3.45, 3.65, 1.0, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_box(slide, "4. Fit only the remaining Kd error to dGDN/dt", 8.9, 4.75, 3.65, 0.92, fill=RGBColor(255, 244, 224), line=ORANGE, size=15, bold=True)
    add_takeaway(slide, "This fit uses the original IBIS-derived Ku/Kd tables only; it does not learn from HSPICE.")
    add_source(slide, "Actual implementation: tools/pybis2spice/pybis2spice/subcircuit.py::two_state_directional_gate_fit")


def build_state_code_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Implementation 4: capacitor-backed hidden states in generated SPICE")
    code = (
        "* GUP: pullup hidden state\n"
        "BGUP GUP 0 I = -{gate_c} * (V(GUPTARGET)-V(GUP)) /\n"
        "+ ((V(GUPTARGET)>V(GUP)) ? 1.28444958348n : 0.363111553053n)\n"
        "CGUP GUP 0 {gate_c} ic=0\n"
        "RGUP GUP 0 1e12\n"
        "\n"
        "* GDN: pulldown hidden state\n"
        "BGDN GDN 0 I = -{gate_c} * (V(GDNTARGET)-V(GDN)) /\n"
        "+ ((V(GDNTARGET)>V(GDN)) ? 0.274962895839n : 0.237281600919n)\n"
        "CGDN GDN 0 {gate_c} ic=1\n"
        "RGDN GDN GDNBASE 1e12"
    )
    add_code_box(slide, code, 0.55, 1.12, 8.55, 4.85, size=12.2)
    add_text(slide, "How to read this", 9.45, 1.22, 2.8, 0.38, size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "B source supplies charging current.",
            "C stores the continuous state voltage.",
            "Different tau is selected for charging and discharging.",
            "Initial condition is low output: GUP=0, GDN=1.",
            "A reverse edge changes the target, not the stored voltage.",
        ],
        9.45,
        1.7,
        3.2,
        3.55,
        size=15.5,
    )
    add_box(slide, "Result: GUP and GDN cannot jump at retrigger.", 9.45, 5.38, 3.1, 0.64, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_takeaway(slide, "The SPICE capacitor is not the physical transistor capacitance; it is a compact state-memory implementation.")
    add_source(slide, "Actual generated subcircuit: cases/short_pulse_1ns_high/.../driver_OutputInput_Typical.sub")


def build_mapping_code_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Implementation 5: choose direction, add residual, emit continuous Ku/Kd")
    code = (
        "BKUGATE_BASE ... V = (V(GUPTARGET)>=V(GUP))\n"
        "+ ? V(KUGATE_ON) : V(KUGATE_OFF)\n"
        "BKDGATE_BASE ... V = (V(GDNTARGET)>=V(GDN))\n"
        "+ ? V(KDGATE_ON) : V(KDGATE_OFF)\n"
        "\n"
        "BGDNRATE ... V = (V(GDNTARGET)-V(GDN))/tau_pd * 1e-9\n"
        "BKDRES ... V = V(KDRES_TABLE)\n"
        "+              + 0.00023869148*V(GDNRATE)\n"
        "BKUGATE KUGATE 0 V = V(KUGATE_BASE)\n"
        "BKDGATE KDGATE 0 V = V(KDGATE_BASE) + V(KDRES)\n"
        "\n"
        "B44 Ku 0 I = -{coeff_c}*(V(KUTARGET)-V(Ku))/coeff_tau\n"
        "Cku Ku 0 {coeff_c} ic=0.00194059\n"
        "B45 Kd 0 I = -{coeff_c}*(V(KDTARGET)-V(Kd))/coeff_tau\n"
        "Ckd Kd 0 {coeff_c} ic=0.99861080"
    )
    add_code_box(slide, code, 0.52, 1.0, 8.25, 5.55, size=10.7)
    add_box(slide, "Direction selector\nTarget above state: ON map\nTarget below state: OFF map", 9.08, 1.18, 3.45, 1.05, fill=LIGHT_GREEN, line=GREEN, size=15, bold=True)
    add_box(slide, "Kd residual\nAdds table error plus a small dGDN/dt term to restore negative undershoot", 9.08, 2.72, 3.45, 1.28, fill=RGBColor(255, 244, 224), line=ORANGE, size=15, bold=True)
    add_box(slide, "Final smoothing\nKu/Kd are capacitor-backed, preventing coefficient jumps", 9.08, 4.5, 3.45, 1.08, fill=RGBColor(234, 244, 255), line=BLUE, size=15, bold=True)
    add_source(slide, "Exact generated syntax shortened only by replacing long PWL point lists with labels.")


def build_reconstruction_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "First validation gate: reproduce the original complete-edge tables")
    path = RESULT_ROOT / "presentation_evidence_figures" / "reconstruction_gate_evidence.png"
    add_picture_contain(slide, path, 0.38, 0.92, 12.55, 5.72)
    add_takeaway(slide, "Directional maps plus the Kd rate residual pass: worst RMSE 0.01994, worst max error 0.04869.")
    add_source(slide, str(path.relative_to(ROOT)))


def add_waveform_slide(prs: Presentation, title: str, image_name: str, bullets: list[str], status_text: str, status_color: RGBColor) -> None:
    slide = add_slide(prs, title)
    path = RESULT_ROOT / "waveform_evidence_report" / "plots" / "directional_residual_vs_native_ibis" / image_name
    add_picture_contain(slide, path, 0.4, 1.0, 9.45, 5.55)
    add_box(slide, status_text, 10.15, 1.18, 2.65, 0.6, fill=PALE_GREEN if status_color == GREEN else LIGHT_RED, line=status_color, size=16, bold=True, color=status_color)
    add_bullets(slide, bullets, 10.08, 2.08, 2.75, 3.62, size=14.5, spacing=8)
    add_takeaway(slide, bullets[-1], y=6.65, fill=PALE_GREEN if status_color == GREEN else LIGHT_RED, color=status_color)
    add_source(slide, str(path.relative_to(ROOT)) + " (cached data only)")


def build_transistor_comparison_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Pad-level check against transistor io_buf.sp")
    low = RESULT_ROOT / "waveform_evidence_report" / "plots" / "directional_residual_vs_hspice_sp" / "short_pulse_2ns_low_directional_residual_vs_hspice_sp.png"
    high = RESULT_ROOT / "waveform_evidence_report" / "plots" / "directional_residual_vs_hspice_sp" / "short_pulse_1ns_high_directional_residual_vs_hspice_sp.png"
    add_text(slide, "2 ns short-low", 0.62, 0.92, 5.95, 0.35, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "1 ns short-high", 6.75, 0.92, 5.95, 0.35, size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_picture_contain(slide, low, 0.45, 1.3, 6.1, 4.45)
    add_picture_contain(slide, high, 6.72, 1.3, 6.1, 4.45)
    add_text(slide, "Transistor SPICE exposes pad voltage only - there is no transistor-level Ku/Kd trace.", 0.75, 5.92, 12.0, 0.35, size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Native IBIS is the coefficient-playback reference; io_buf.sp is an independent pad-level circuit reference.")
    add_source(slide, "Cached HSPICE transistor data; no simulations were rerun for this presentation.")


def build_results_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Measured status: strongest structure, still an experimental model")
    add_metric_card(slide, "Offline reconstruction", "RMSE 0.01994\nMax 0.04869", "PASS", 0.55, 1.18, 2.95, GREEN)
    add_metric_card(slide, "Normal long pulse", "Pad 18.739 mV\nKu .01886 | Kd .01110", "WARN: legacy is closer", 3.72, 1.18, 2.95, ORANGE)
    add_metric_card(slide, "2 ns short-low", "Pad 23.271 mV\nKu .02141 | Kd .01319", "GOOD", 6.89, 1.18, 2.95, GREEN)
    add_metric_card(slide, "1 ns short-high", "Pad 21.914 mV\nKu .02161 | Kd .48754", "OPEN: Kd recovery", 10.06, 1.18, 2.7, RED, value_size=13)

    add_text(slide, "What is now proven", 0.72, 3.05, 3.3, 0.38, size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "Two continuous states can reproduce complete-edge Ku/Kd when direction-specific maps are used.",
            "The Kd rate residual restores the negative excursion that static maps miss.",
            "Short-low behavior demonstrates coefficient-correct improvement in at least one interrupted quadrant.",
        ],
        0.76,
        3.48,
        5.75,
        2.1,
        size=16,
    )
    add_text(slide, "What is not yet proven", 6.85, 3.05, 3.6, 0.38, size=20, bold=True, color=RED)
    add_bullets(
        slide,
        [
            "The full model is not a drop-in replacement for legacy pybis on normal edges.",
            "Short-high Kd onset/hold/recovery is still wrong even when pad voltage looks small.",
            "Native IBIS and transistor io_buf.sp differ, so coefficient agreement is playback agreement, not transistor truth.",
        ],
        6.88,
        3.48,
        5.65,
        2.1,
        size=16,
    )
    add_takeaway(slide, "Conclusion: two_state_directional_residual is the best structural method so far, not the production default.")
    add_source(slide, "Metrics: results/io_buf_two_state_gate_model_2026-06-30/candidate_metrics.csv and README.md")


def build_limits_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Why the remaining short-high Kd problem is difficult")
    add_box(slide, "Static state shape\nmostly solved", 0.7, 1.32, 2.45, 1.0, fill=LIGHT_GREEN, line=GREEN, size=17, bold=True)
    add_arrow(slide, 3.15, 1.82, 4.0, 1.82)
    add_box(slide, "Recovery onset / hold\nstill width-sensitive", 4.0, 1.32, 2.8, 1.0, fill=LIGHT_RED, line=RED, size=17, bold=True)
    add_arrow(slide, 6.8, 1.82, 7.65, 1.82)
    add_box(slide, "Present GDN@reverse\ndoes not encode enough history", 7.65, 1.32, 3.05, 1.0, fill=LIGHT_RED, line=RED, size=16, bold=True)
    add_arrow(slide, 10.7, 1.82, 11.45, 1.82)
    add_box(slide, "Need richer\ncausal state", 11.45, 1.32, 1.3, 1.0, fill=LIGHT_GRAY, line=GRAY, size=15, bold=True)

    add_text(slide, "Evidence already collected", 0.72, 3.05, 3.2, 0.38, size=20, bold=True, color=GREEN)
    add_bullets(
        slide,
        [
            "Main-slope Kd tau changes only about 1.086x across widths: changing tau alone is not enough.",
            "A fixed mean recovery delay helps some widths but is early or late for others.",
            "GDN@reverse collapses 500 ps and 1 ns pulses to nearly the same state.",
            "A simple command-age hold law failed the 1.5 ns held-out validation by 49.4 ps.",
        ],
        0.78,
        3.48,
        6.0,
        2.25,
        size=15.5,
    )
    add_text(slide, "Engineering implication", 7.18, 3.05, 3.2, 0.38, size=20, bold=True, color=RED)
    add_box(slide, "The next state must track pending predriver history or recovery phase - not just present GDN voltage.", 7.25, 3.65, 5.0, 1.35, fill=LIGHT_RED, line=RED, size=18, bold=True)
    add_box(slide, "Any new candidate must preserve long-pulse legacy behavior and pass Ku/Kd before pad voltage.", 7.25, 5.2, 5.0, 0.72, fill=PALE_GREEN, line=GREEN, size=15.5, bold=True)
    add_takeaway(slide, "The current result identifies the missing state variable; it does not justify per-pulse tuning.")
    add_source(slide, "Diagnostics: kd_recovery_diagnostics/effective_tau, gdn_hold_time, command_age_hold")


def build_next_steps_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Next steps toward a trustworthy default")
    steps = [
        ("1", "Preserve normal edges", "Use legacy-equivalent output outside a proven interrupted-transition window."),
        ("2", "Add causal recovery memory", "Track pending pulldown-off/recovery phase separately from GDN amplitude."),
        ("3", "Validate both directions", "Short-high and short-low widths, plus repeated toggles and load changes."),
        ("4", "Keep coefficient-first gates", "Ku and Kd must improve together; pad-only agreement cannot pass."),
        ("5", "Broaden model coverage", "Repeat on additional IBIS buffers before changing the pybis default."),
    ]
    y = 1.08
    for number, head, detail in steps:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.78), Inches(y), Inches(0.52), Inches(0.52))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN
        circle.line.color.rgb = GREEN
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = number
        p.alignment = PP_ALIGN.CENTER
        p.font.name = BODY_FONT
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        add_text(slide, head, 1.52, y - 0.02, 3.6, 0.42, size=16.5, bold=True, color=GREEN)
        add_text(slide, detail, 5.32, y - 0.03, 7.05, 0.58, size=15.2)
        y += 1.05
    add_takeaway(slide, "Keep the two-state directional architecture; improve only the missing recovery-memory mechanism.")
    add_source(slide, "Recommended direction based on cached June 30 - July 10 studies.")


def build_closing_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "Takeaway")
    add_box(slide, "Problem", 0.82, 1.25, 3.55, 0.62, fill=LIGHT_RED, line=RED, size=19, bold=True, color=RED)
    add_text(slide, "Legacy replay restarts from a settled assumption when a pulse reverses early.", 0.82, 2.05, 3.55, 1.05, size=18, align=PP_ALIGN.CENTER)
    add_box(slide, "Solution", 4.9, 1.25, 3.55, 0.62, fill=LIGHT_GREEN, line=GREEN, size=19, bold=True, color=GREEN)
    add_text(slide, "Store continuous pullup/pulldown gate states, map by direction, and restore Kd dynamics with a rate residual.", 4.9, 2.05, 3.55, 1.25, size=18, align=PP_ALIGN.CENTER)
    add_box(slide, "Current status", 8.98, 1.25, 3.55, 0.62, fill=RGBColor(255, 244, 224), line=ORANGE, size=19, bold=True, color=ORANGE)
    add_text(slide, "Best structural method so far. Reconstruction passes; short-high Kd recovery remains open.", 8.98, 2.05, 3.55, 1.05, size=18, align=PP_ALIGN.CENTER)

    add_box(slide, "two_state_directional_residual", 2.12, 4.18, 9.1, 0.88, fill=GREEN, line=GREEN, size=27, bold=True, color=WHITE)
    add_text(slide, "A credible architecture and a measured research result - not yet the production default.", 2.05, 5.35, 9.25, 0.62, size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_takeaway(slide, "Questions")
    add_source(slide, "Simon Hwang | UMR EMC Laboratory | 7/14/2026")


PRESENTER_NOTES = {
    "Short-Pulse IBIS Simulation in Ngspice": """
Open by stating the practical goal: reproduce IBIS buffer behavior in ngspice when the input reverses before the previous output transition has settled. The project is not trying to replace the IBIS standard or claim transistor-level accuracy from Ku and Kd. It is testing whether a compact, continuous-state implementation can avoid the obvious restart error of legacy table replay.

The talk proceeds in three layers. First, introduce Ku and Kd and show why a short pulse is different from a normal edge. Second, derive the two hidden-state model from the original complete-edge coefficient tables. Third, show the generated SPICE syntax and the cached validation evidence. Emphasize that HSPICE is an audit reference only; it does not provide fitting data.
""",
    "What are Ku and Kd?": """
Explain the output buffer as two controlled networks connected to one pad. Ku scales the pullup I-V table and Kd scales the pulldown I-V table. A coefficient near one means that network contributes close to its fully enabled table; a coefficient near zero means it is mostly disabled. Small overshoot or undershoot outside zero to one can occur in the extracted tables, so these are effective coefficients rather than literal transistor gate voltages.

For a rising output, Ku normally increases while Kd decreases. For a falling output, Ku decreases while Kd increases. In a long pulse both processes eventually settle. In a short pulse the reverse edge arrives while one or both coefficients are still moving. The central modeling problem is therefore history: the second edge must continue from an unfinished internal condition.
""",
    "\"Replay Problem\" for Short Pulse": """
Use this slide to define the legacy failure. The normal pybis InputDriven path detects an edge and indexes a complete rising or falling Ku/Kd table by elapsed time since that edge. That works when every edge begins from the settled endpoint assumed by the table. It does not work when a reverse edge arrives early.

At the reverse edge, the old elapsed timer is restarted and the opposite complete-edge table begins from its own t=0 endpoint. That endpoint represents a settled logic state, not the partial state actually reached. The figure illustrates the resulting coefficient jump or full-strength replay. The error is structural: improving the timestep or output load cannot repair the wrong initial condition.
""",
    "Why \"Value-matched\" Method Fails": """
This slide summarizes the value-matched baseline. At the reverse edge, that method samples the current Ku and Kd and searches the opposite tables for times that produce similar values. It then tries to restart the opposite replay from that inferred table time.

The key issue is that Ku and Kd usually imply different start times. Their rising and falling trajectories are asymmetric, so one shared replay time cannot simultaneously match both sampled values. Separate starts improve one coefficient but break their coordinated timing, while a balanced start is a compromise that is not a physical state. The earlier v1 implementation also had a timer bug, but the corrected v2 experiment still showed table-start ambiguity. This motivates a model with independent continuous internal state rather than a single retimed table cursor.
""",
    "Controlled experiment: same stimulus, three references": """
Walk through the testbench from left to right. The command is a 3.3 V digital pulse with 1 ps edges. The output drives a direct 50 ohm parallel 2 pF load. The same stimulus and load are used for each flow.

HSPICE native IBIS exposes pad voltage plus internal Ku and Kd and is the coefficient-playback reference. The HSPICE transistor io_buf.sp run exposes pad voltage only and is a separate circuit-level reference. The ngspice candidate exposes pad, Ku, Kd, and hidden state diagnostics. The long pulse is a regression guard. Short-high and short-low pulses test interrupted transitions. Cached HSPICE results are reused whenever the deck and model hashes match, and no HSPICE data is used to derive parameters.
""",
    "Proposed solution: remember two hidden gate states": """
Introduce the architecture at a high level before discussing equations. Input edges first pass through four direction-specific delays. The delayed commands drive two continuous states: GUP for pullup progress and GDN for pulldown progress. Direction-specific PWL maps convert those states into Ku and Kd. A small Kd residual restores dynamic detail that the static state map misses.

At a reverse edge, only the targets change. GUP and GDN retain their present values and reverse continuously. This is the central difference from table replay. The method name two_state_directional_residual describes the three structural pieces: two stored states, different on/off mappings, and a Kd residual correction.
""",
    "How a capacitor stores unfinished switching history": """
The capacitor is a numerical state-storage device in the generated SPICE, not a claim that it equals a particular transistor capacitance. A behavioral current source applies C times (target minus state) divided by tau. Because I = C dV/dt, the capacitor voltage obeys dG/dt = (target - G)/tau.

If a reverse edge occurs when GUP is 0.31 and GDN is 0.42, the targets change but the capacitor voltages cannot jump. The next trajectory begins at exactly 0.31 and 0.42. Separate charging and discharging taus are selected according to whether the target is above or below the current state. Large leak resistors provide a DC path without materially changing the transient.
""",
    "Direction matters: four processes, four fitted dynamics": """
There are four independently fitted processes. On a rising input, the pullup turns on with 1.302 ns delay and 1.284 ns tau, while the pulldown turns off with 1.443 ns delay and 0.237 ns tau. On a falling input, the pullup turns off with 0.440 ns delay and 0.363 ns tau, while the pulldown turns on with 2.480 ns delay and 0.275 ns tau.

These labels and numbers match gate_fit_summary.csv and the generated subcircuit. The runtime map choice is also directional: target above state selects the on map, and target below state selects the off map. A single shared map failed because the same normalized state can correspond to different coefficient values on the on and off trajectories.
""",
    "Step 0: the model learns from four complete-edge traces": """
Define kr as the rising-input coefficient table and kf as the falling-input coefficient table. Each row contains time, Ku, and Kd. Therefore the available training traces are Ku-on and Kd-off from kr, plus Ku-off and Kd-on from kf.

The endpoint estimator averages redundant settled values from both tables. Ku_off averages the beginning of the rising table and the end of the falling table. Ku_on averages the end of rising and beginning of falling. Kd uses the analogous pairings. The io_buf values shown on the slide are the exact values recorded in gate_fit_summary.csv. This is the full fitting input. HSPICE waveforms are deliberately absent.
""",
    "Step 1: normalize each coefficient into transition progress": """
Normalization removes the coefficient's endpoint scale. Progress is (K - K_start) divided by (K_end - K_start). If K decreases, the denominator is negative, so progress still increases from zero to one. This allows one timing extractor to work for all four processes.

The implementation clips progress to -0.5 through 1.5, not zero through one. That wider diagnostic range preserves real table overshoot and undershoot for crossing detection and later residual analysis. crossing_time_ns sorts the samples and linearly interpolates the first crossing of each requested level. It does not assume uniformly spaced samples.
""",
    "Step 2: derive delay and tau from 5%, 63.2%, and 90% crossings": """
Explain the first-order template. Before the delay, the hidden state is stationary. After the delay it follows 1 - exp(-(t-delay)/tau) for an increasing transition. A first-order response reaches 63.2 percent after one tau and 90 percent after ln(10), or 2.3026, taus.

The code sets delay to t05. It obtains one tau estimate from t63 minus delay and another from (t90 minus delay) divided by ln(10). It chooses the largest of those and a 20 ps floor. Choosing the larger estimate protects against a fast initial segment followed by a slower tail. This is a deterministic crossing-based estimate, not a nonlinear least-squares optimization of the whole trace.
""",
    "Step 3: reconstruct GUP and GDN from the fitted timing": """
For every original table time, gate_response computes x = max(0, t-delay), progress = 1-exp(-x/tau), and G = G_start + (G_end-G_start) times progress. Increasing and decreasing states use the same equation through their endpoints.

GUP starts at zero for a settled low output and rises toward one when the pullup is commanded on. GDN starts at one and falls toward zero when the pulldown is commanded off. The falling input uses the opposite endpoints. GUP and GDN are latent normalized progress variables. They are not assumed equal to Ku and Kd; the next stage learns the nonlinear mapping from each G state to the measured coefficient.
""",
    "Step 4: turn hidden state and coefficient samples into four PWL maps": """
At each table sample time, the code now has both the fitted hidden state G(t) and the original coefficient K(t). It pairs them, sorts the pairs by G, averages duplicate G locations, and interpolates onto 81 evenly spaced state points from zero to one. Endpoints are set to the inferred settled coefficients.

The operation is repeated independently for Ku-on, Ku-off, Kd-off, and Kd-on. Unlike the original single-map helper, the directional helper does not force the coefficient values to be monotonic. This is intentional because forcing monotonicity erased the negative Kd excursion. At runtime, the sign of target minus state selects the correct on or off map.
""",
    "Step 5: recover the Kd detail that a static state map misses": """
First reconstruct Kd using only the direction-specific state map. Subtract that reconstruction from the original Kd table to obtain the residual. Separately compute dGDN/dt from the fitted first-order state. Offline this is the analytic exponential derivative; in SPICE it is evaluated causally as (target-GDN)/tau.

The scalar rate gain is the least-squares projection of residual onto rate: dot(rate,residual) divided by dot(rate,rate). Its unit is ns because rate is state per ns. For io_buf the gain is 0.0002386915 ns. Runtime Kd adds three pieces: directional base map, an elapsed-edge-time residual table, and gain times dGDN/dt. Be explicit that the main state path is continuous, while the small residual-table term still retains edge-time indexing. That is a known architectural nuance, not hidden from the audience.
""",
    "Reading the fit evidence: the on and off maps are genuinely different": """
Read the axes before interpreting the curves. The horizontal axis is normalized hidden state, not time. The vertical axis is Ku or Kd. Blue and orange are the two directions for the same network.

The separation between the curves demonstrates directional hysteresis in the extracted representation: the same GUP or GDN value maps to different effective coefficients depending on whether the network is turning on or off. A single map would average across those branches and create reconstruction error. The plotted curves are the actual 81-point PWL data generated from the complete-edge tables.
""",
    "Reading the fit evidence: rate and residual explain the Kd undershoot": """
The left plot shows dGDN/dt versus original table time. It is negative when the pulldown state is turning off and positive when it turns on. The right plot shows the residual after subtracting the direction-specific PWL reconstruction from the original Kd table.

The strong negative region is the important feature. The earlier identity and static PWL models could not reproduce it and therefore failed the normal reconstruction gate. The residual table plus the small fitted rate term restores that feature well enough for the directional-residual variant to pass the offline gate. Passing this gate is necessary, but it does not by itself prove interrupted-pulse behavior.
""",
    "End-to-end implementation flow": """
Follow the artifacts from left to right. pybis extracts complete-edge Ku/Kd coefficient tables from io_buf.ibs. The directional fit converts those tables into endpoints, four delays, four taus, four PWL maps, two Kd residual traces, and one rate gain. The SPICE generator emits behavioral sources, transport delays, state capacitors, map lookup sources, and final coefficient nodes.

The resulting subcircuit is inserted into an ngspice transient deck. Independent cached HSPICE native-IBIS and transistor-pad runs are loaded only during reporting. The normal generation path therefore requires the IBIS file and ngspice, not HSPICE.
""",
    "Implementation 1: endpoint and timing extraction in Python": """
This slide is close to the actual Python. The endpoint expressions show exactly which samples are averaged. coefficient_progress performs the normalization described earlier. crossing_time_ns finds the first interpolated crossing after sorting valid samples.

The timing function applies the 5, 63.2, and 90 percent recipe. The 20 ps minimum prevents an unrealistically zero time constant from producing a numerically stiff or singular state equation. Four calls produce PU-on, PU-off, PD-on, and PD-off timing. The code uses nan-aware operations and finite fallbacks so a malformed endpoint does not silently propagate NaN into the netlist.
""",
    "Implementation 2: build directional PWL maps and fit the rate gain": """
The first two lines reconstruct GDN during the rising and falling complete-edge tables using their respective delays, taus, and endpoints. directional_gate_transfer_curve creates the two Kd maps by pairing each GDN sample with the original Kd sample.

np.interp evaluates the map back at the reconstructed state to create Kd_base. Subtracting from the original table produces the residual. The dot-product expression is a one-parameter least-squares fit of residual against dGDN/dt. Ku follows the same map-construction process but does not receive this rate-residual term in the current best structural variant.
""",
    "Implementation 3: each input edge launches four delayed gate commands": """
RISEEDGE and FALLEDGE are short event pulses generated by the input detector. Each event is copied through a transport delay implemented by a lossless 50 ohm T-line with a matched resistor. The four delayed nodes are PUONP, PUOFFP, PDOFFP, and PDONP.

Those pulses add or remove charge from GUPCMD and GDNCMD command capacitors. Clamped behavioral sources convert the command levels into GUPTARGET and GDNTARGET in the zero-to-one range. Importantly, these are target states. The actual GUP and GDN capacitors are separate and evolve continuously. Pending delayed commands are not represented by simply restarting a Ku/Kd table.
""",
    "Implementation 4: capacitor-backed hidden states in generated SPICE": """
Read the GUP branch first. BGUP supplies current proportional to target minus state, divided by the on or off tau. CGUP stores the state voltage and starts at zero. The large resistor gives the node a DC reference. Combining I=C dV/dt with the source equation gives dGUP/dt=(target-GUP)/tau.

GDN uses the same structure but starts at one. When GDNTARGET is above GDN, the pulldown-on tau is selected; when it is below, the pulldown-off tau is selected. A reverse command changes only the target and selected tau. The capacitor state itself cannot jump. gate_c cancels algebraically from the intended differential equation; its 1 pF value is a convenient numerical scaling.
""",
    "Implementation 5: choose direction, add residual, emit continuous Ku/Kd": """
The first two behavioral sources choose an on or off PWL value by comparing target with current state. BGDNRATE computes dGDN/dt in per-ns units; multiplying the SI per-second expression by 1e-9 performs that unit conversion. BKDRES adds the residual table and the fitted rate contribution to Kd_base.

KUGATE and KDGATE are target coefficients. Final Ku and Kd each have another small 5 ps capacitor-backed smoothing stage. This prevents ideal algebraic jumps when a map branch changes. These final nodes scale the IBIS pullup and pulldown I-V sources. The long PWL point lists are omitted on the slide, but the syntax and numerical constants match the generated subcircuit.
""",
    "First validation gate: reproduce the original complete-edge tables": """
This is the legitimacy gate applied before interpreting short pulses. The single-map model has worst max error about 0.205 and fails. Direction-specific maps reduce that to about 0.074 but still miss the required threshold. Adding the Kd residual reduces worst RMSE to 0.01994 and worst max error to 0.04869, which passes the predefined offline gate.

The black curves are the original extracted tables and red curves are model reconstructions. This is not transient simulation output. It asks whether the compact representation can reproduce the normal complete-edge coefficient data from which it was derived.
""",
    "Normal transition: structure is close, but legacy remains better": """
The long-pulse case protects behavior that legacy pybis already handles well. The directional-residual Ku and Kd shapes are close to native IBIS, but its pad RMSE is 18.739 mV versus 5.289 mV for legacy. Therefore this case is a warning, not a win.

Point out that passing the offline reconstruction gate did not guarantee identical closed-loop transient behavior. Small coefficient timing differences interact with the nonlinear IBIS I-V tables and load. Any future promotion must preserve or explicitly hybridize the normal path rather than accepting this regression.
""",
    "Short-low success: right for the right reasons": """
This is the strongest success case in the study. For the 2 ns short-low pulse, pad voltage, Ku, and Kd improve together relative to the coefficient reference. The directional-residual result has pad RMSE 23.271 mV, Ku RMSE 0.02141, and Kd RMSE 0.01319.

Stress the interpretation rule: a pad waveform alone is not enough. This case is useful because the internal coefficients and output consequence agree at the same time. It demonstrates that the two-state directional structure can solve at least one interrupted-transition quadrant without a table restart.
""",
    "Short-high remains open: pad alone hides the Kd error": """
For the 1 ns short-high pulse, Ku remains partial and is close to native IBIS. The pad signal is also small, so its absolute RMSE of 21.914 mV looks attractive. However, Kd recovery is dramatically different and Kd RMSE is 0.48754.

This is the clearest example of the coefficient-first rule. The load and pullup behavior make the small pad excursion relatively insensitive to the wrong pulldown coefficient, so pad agreement can hide a real model error. The current architecture still lacks the correct short-high Kd recovery staging or history variable.
""",
    "Pad-level check against transistor io_buf.sp": """
The transistor-level io_buf.sp run is a pad-only reference. It does not expose Ku or Kd because those are IBIS playback coefficients, not physical transistor observables. The plots therefore compare only pad voltage against the directional-residual model.

The native IBIS and transistor references can differ substantially, including about 525 mV on the long-pulse pad in this study. That changes how claims must be phrased. Ku/Kd agreement measures consistency with native-IBIS playback. Transistor-pad agreement is a separate check of circuit-level output behavior. Neither should be silently substituted for the other.
""",
    "Measured status: strongest structure, still an experimental model": """
Summarize the four gates. Offline complete-edge reconstruction passes. Normal long-pulse transient remains a warning because legacy is better. The 2 ns short-low case is good in pad and coefficients. The 1 ns short-high case remains open because Kd recovery is wrong.

The proper conclusion is therefore structural progress, not production readiness. The method has demonstrated that continuous GUP/GDN states, direction-specific maps, and a rate residual are useful. It has not demonstrated a general replacement for legacy InputDriven across all pulse directions and widths.
""",
    "Why the remaining short-high Kd problem is difficult": """
The static map shape is mostly solved, but recovery onset and hold behavior vary with interruption depth. A fixed recovery delay improves some widths and hurts others. The main-slope tau changes only about 1.086 times, so merely changing tau does not explain the missing early trajectory.

The current GDN value at the reverse edge is also insufficient: it collapses the 500 ps and 1 ns cases to nearly the same state. A pulse-width hold law fit the training points but failed the held-out 1.5 ns case by 49.4 ps against a 30 ps gate. The evidence points to a missing causal recovery-phase or pending-command state, not a need for per-case tuning.
""",
    "Next steps toward a trustworthy default": """
The next candidate should retain the current architecture and add only the missing recovery memory. Outside a proven interrupted window, preserve legacy-equivalent normal behavior. During interruption, track a causal state representing pending pulldown-off or recovery phase in addition to GDN amplitude.

Validation must cover short-high, short-low, repeated toggles, multiple widths, and load changes. Ku and Kd remain primary gates, with pad as the physical consequence. Finally, repeat the method on additional IBIS buffers before changing the pybis default; a solution tuned only to io_buf is not a general algorithm.
""",
    "Takeaway": """
Close with the three-part message. The problem is a restart assumption: legacy elapsed-time replay treats a reverse edge as if the previous transition had fully settled. The proposed solution stores two continuous gate-progress states, uses separate on/off maps, and adds a small Kd dynamic correction.

The method is the best structural result so far because it passes complete-edge reconstruction and gives a coefficient-correct short-low result. It is not the default because normal-edge pad behavior regresses and short-high Kd recovery remains wrong. The next work is targeted: add a causal recovery-memory state while preserving the validated pieces and the long-pulse guard.
""",
}


def apply_presenter_notes(prs: Presentation) -> None:
    missing = []
    for slide in prs.slides:
        title = slide.shapes.title.text.strip() if slide.shapes.title is not None else ""
        note = PRESENTER_NOTES.get(title)
        if note is None:
            missing.append(title or "<untitled>")
            continue
        add_notes(slide, note)
    if missing:
        raise ValueError(f"Missing presenter notes for slides: {missing}")


def update_title_slide(prs: Presentation) -> None:
    slide = prs.slides[0]
    title = slide.shapes.title
    title.text = "Short-Pulse IBIS Simulation in Ngspice"
    p = title.text_frame.paragraphs[0]
    p.font.name = TITLE_FONT
    p.font.size = Pt(37)
    p.font.bold = True
    subtitle = slide.placeholders[1]
    subtitle.text = "Simon Hwang\nAdvisor: Dr. Chulsoon Hwang\nDr. Zhiping Yang\n7/14/2026"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = TITLE_FONT
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = BLACK


def build_deck(template: Path, output: Path) -> None:
    prs = Presentation(str(template))
    update_title_slide(prs)
    build_primer(prs)
    build_setup_slide(prs)
    build_architecture_slide(prs)
    build_state_memory_slide(prs)
    build_direction_slide(prs)
    build_training_data_slide(prs)
    build_progress_slide(prs)
    build_tau_derivation_slide(prs)
    build_hidden_state_reconstruction_slide(prs)
    build_pwl_derivation_slide(prs)
    build_rate_residual_derivation_slide(prs)
    build_map_diagnostic_slide(prs)
    build_residual_diagnostic_slide(prs)
    build_software_flow_slide(prs)
    build_timing_code_detail_slide(prs)
    build_map_fit_code_detail_slide(prs)
    build_command_target_slide(prs)
    build_state_code_slide(prs)
    build_mapping_code_slide(prs)
    build_reconstruction_slide(prs)
    add_waveform_slide(
        prs,
        "Normal transition: structure is close, but legacy remains better",
        "edge_1ps_base_50r_2pf_directional_residual_vs_native_ibis.png",
        [
            "Ku/Kd shapes track native IBIS closely.",
            "Pad RMSE is 18.739 mV versus 5.289 mV for legacy.",
            "This is a regression guard, not a claimed win.",
            "Normal behavior must be preserved before promotion.",
        ],
        "NORMAL EDGE: WARN",
        ORANGE,
    )
    add_waveform_slide(
        prs,
        "Short-low success: right for the right reasons",
        "short_pulse_2ns_low_directional_residual_vs_native_ibis.png",
        [
            "Pad, Ku, and Kd improve together.",
            "Pad RMSE: 23.271 mV.",
            "Ku RMSE: 0.02141; Kd RMSE: 0.01319.",
            "This is the strongest interrupted-transition result.",
        ],
        "2 NS SHORT-LOW: GOOD",
        GREEN,
    )
    add_waveform_slide(
        prs,
        "Short-high remains open: pad alone hides the Kd error",
        "short_pulse_1ns_high_directional_residual_vs_native_ibis.png",
        [
            "Ku stays partial and is close to native IBIS.",
            "Absolute pad error is small because output swing is tiny.",
            "Kd recovery is late and strongly different.",
            "Kd RMSE 0.48754 prevents a coefficient-correct pass.",
        ],
        "1 NS SHORT-HIGH: OPEN",
        RED,
    )
    build_transistor_comparison_slide(prs)
    build_results_slide(prs)
    build_limits_slide(prs)
    build_next_steps_slide(prs)
    build_closing_slide(prs)

    apply_presenter_notes(prs)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the two-state directional-residual presentation.")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    build_deck(args.template, args.output)
    print(args.output)


if __name__ == "__main__":
    main()
